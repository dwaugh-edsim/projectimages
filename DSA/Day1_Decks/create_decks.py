import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

out_dir = r"Z:\simroom\DSA\Day1_Decks"
os.makedirs(out_dir, exist_ok=True)

# Color Palette (Anti-AI Slop: Deep Slate Navy, Warm Sand / Bone, Gold/Amber Accent)
BG_COLOR = RGBColor(18, 24, 38)        # #121826 Deep Slate Navy
CARD_BG = RGBColor(28, 36, 56)         # #1C2438 Dark Slate Card
TEXT_WHITE = RGBColor(240, 244, 248)   # #F0F4F8 Crisp Off-White
TEXT_MUTED = RGBColor(156, 169, 186)   # #9CA9BA Muted Slate
ACCENT_GOLD = RGBColor(230, 168, 34)   # #E6A822 Warm Gold
ACCENT_TEAL = RGBColor(45, 170, 185)   # #2DAAB9 Teal Accent

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_hl9_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1: Title & Agenda ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_COLOR)

    # Title Box
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "HEALTHY LIVING 9"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_GOLD
    p1.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = "The Attention Economy & Personal Telemetry"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.font.name = "Segoe UI"

    # Agenda Card
    card = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.0), Inches(11.333), Inches(3.6))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_GOLD
    card.line.width = Pt(1.5)

    ctf = card.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = "TODAY'S CLASS RUN OF SHOW"
    cp.font.size = Pt(18)
    cp.font.bold = True
    cp.font.color.rgb = ACCENT_GOLD
    cp.font.name = "Segoe UI"

    agenda_items = [
        ("1. Core Concept:", "Health is not just nutrition and cardio—it begins with what captures your consciousness."),
        ("2. Live Telemetry Pull:", "Extract 3 un-fakeable data points directly from your pocket device settings right now."),
        ("3. Partner Interrogation:", "3-minute verbal peer interviews analyzing triggers, dopamine hooks, and recovery."),
        ("4. Class Attention Heatmap:", "Live anonymous board tally mapping our room's daily baseline.")
    ]

    for num, desc in agenda_items:
        p = ctf.add_paragraph()
        p.space_before = Pt(10)
        run1 = p.add_run()
        run1.text = f"{num} "
        run1.font.bold = True
        run1.font.size = Pt(15)
        run1.font.color.rgb = TEXT_WHITE
        run1.font.name = "Segoe UI"
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(15)
        run2.font.color.rgb = TEXT_MUTED
        run2.font.name = "Segoe UI"

    # --- SLIDE 2: The Core Concept ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BG_COLOR)

    tb2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.2))
    tf2 = tb2.text_frame
    p = tf2.paragraphs[0]
    p.text = "THE ATTENTION ECONOMY"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.font.name = "Segoe UI"
    p_sub = tf2.add_paragraph()
    p_sub.text = "You Are Not Just the User. Your Consciousness is the Harvest."
    p_sub.font.size = Pt(32)
    p_sub.font.bold = True
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.font.name = "Segoe UI"

    # 3 Column Cards
    col_w = Inches(3.5)
    gap = Inches(0.4)
    start_x = Inches(1.0)
    y_pos = Inches(2.4)
    h_pos = Inches(4.2)

    cols = [
        ("Cognitive Scarcity", "Human focus is a finite biological resource. Every push notification is an engineered interruption designed to extract that resource for platform revenue."),
        ("The 5 Dimensions", "Screen habits directly impact all 5 health dimensions: Sleep architecture (Physical), anxiety spikes (Mental), validation loops (Emotional), isolation (Social), and reflection (Spiritual)."),
        ("Zero Fluff / Real Data", "No generic essays. No AI prompts. Today we look at your actual hardware telemetry to understand the exact mechanics of your daily focus.")
    ]

    for i, (col_title, col_text) in enumerate(cols):
        cx = start_x + i * (col_w + gap)
        c_shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y_pos, col_w, h_pos)
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = CARD_BG
        c_shape.line.color.rgb = ACCENT_TEAL if i == 1 else ACCENT_GOLD
        c_shape.line.width = Pt(1)

        ctf = c_shape.text_frame
        ctf.word_wrap = True
        ctp = ctf.paragraphs[0]
        ctp.text = col_title
        ctp.font.size = Pt(18)
        ctp.font.bold = True
        ctp.font.color.rgb = ACCENT_GOLD if i != 1 else ACCENT_TEAL
        ctp.font.name = "Segoe UI"

        ctp_desc = ctf.add_paragraph()
        ctp_desc.space_before = Pt(14)
        ctp_desc.text = col_text
        ctp_desc.font.size = Pt(14)
        ctp_desc.font.color.rgb = TEXT_MUTED
        ctp_desc.font.name = "Segoe UI"

    # --- SLIDE 3: Telemetry Instructions (STAYS ON SCREEN) ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, BG_COLOR)

    tb3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.2))
    tf3 = tb3.text_frame
    p = tf3.paragraphs[0]
    p.text = "STUDENT ACTION PROTOCOL"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.font.name = "Segoe UI"
    p_main = tf3.add_paragraph()
    p_main.text = "Step 1: Phone Telemetry  |  Step 2: Partner Interrogation"
    p_main.font.size = Pt(30)
    p_main.font.bold = True
    p_main.font.color.rgb = TEXT_WHITE
    p_main.font.name = "Segoe UI"

    # Left Card: Telemetry Pull
    card_l = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(5.4), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = ACCENT_TEAL
    card_l.line.width = Pt(1.5)

    ltf = card_l.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = "STEP 1: UNLOCK YOUR STATS (2 MINS)"
    lp.font.size = Pt(16)
    lp.font.bold = True
    lp.font.color.rgb = ACCENT_TEAL
    lp.font.name = "Segoe UI"

    l_items = [
        "Open Settings > Screen Time (iOS) or Digital Wellbeing (Android).",
        "Record 3 exact numbers on a sticky note or in your notes:",
        "• Number 1: Total Screen Time yesterday (Hours & Mins)",
        "• Number 2: Total Pickups (Times unlocked yesterday)",
        "• Number 3: First App Opened within 15 minutes of waking up"
    ]
    for item in l_items:
        p = ltf.add_paragraph()
        p.space_before = Pt(8)
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE if "Number" in item else TEXT_MUTED
        p.font.name = "Segoe UI"

    # Right Card: Partner Interrogation
    card_r = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = ACCENT_GOLD
    card_r.line.width = Pt(1.5)

    rtf = card_r.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = "STEP 2: THE 3-MINUTE INTERROGATION"
    rp.font.size = Pt(16)
    rp.font.bold = True
    rp.font.color.rgb = ACCENT_GOLD
    rp.font.name = "Segoe UI"

    r_items = [
        "Turn to your assigned partner. Partner A interviews B (3 min), then switch (3 min):",
        "1. 'What is your #1 pickup trigger, and what emotion drives it (boredom, habit, FOMO)?'",
        "2. 'If you regained 2 hours of that time daily, what physical or creative project would you build instead?'",
        "3. 'Rate your sleep quality this week from 1-10. How does your bedtime phone routine affect it?'",
        "4. Wrap-up: Place an anonymous tally on the front board for your pickup range."
    ]
    for item in r_items:
        p = rtf.add_paragraph()
        p.space_before = Pt(8)
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE if item.startswith(("1.", "2.", "3.", "4.")) else TEXT_MUTED
        p.font.name = "Segoe UI"

    hl9_path = os.path.join(out_dir, "HL9_Day1_Attention_Economy.pptx")
    prs.save(hl9_path)
    print(f"Saved: {hl9_path}")

def create_cit9_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1: Title & Agenda ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_COLOR)

    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "CITIZENSHIP 9"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_GOLD
    p1.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = "Power, Rights & The Classroom Constitution"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.font.name = "Segoe UI"

    # Agenda Card
    card = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.0), Inches(11.333), Inches(3.6))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_GOLD
    card.line.width = Pt(1.5)

    ctf = card.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = "TODAY'S CLASS RUN OF SHOW"
    cp.font.size = Pt(18)
    cp.font.bold = True
    cp.font.color.rgb = ACCENT_GOLD
    cp.font.name = "Segoe UI"

    agenda_items = [
        ("1. The Nature of Citizenship:", "Citizenship is not a spectator sport. You are active co-governors of Room [X]."),
        ("2. The 4-Corners Stand-Up:", "Physical movement debate testing real classroom governance scenarios."),
        ("3. The 60-Second Civic Pitch:", "Live partner cross-examination establishing your project role and non-negotiables."),
        ("4. Drafting the Constitution:", "Locking in our room's 3 core operational covenants for 2026-2027.")
    ]

    for num, desc in agenda_items:
        p = ctf.add_paragraph()
        p.space_before = Pt(10)
        run1 = p.add_run()
        run1.text = f"{num} "
        run1.font.bold = True
        run1.font.size = Pt(15)
        run1.font.color.rgb = TEXT_WHITE
        run1.font.name = "Segoe UI"
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(15)
        run2.font.color.rgb = TEXT_MUTED
        run2.font.name = "Segoe UI"

    # --- SLIDE 2: The Core Concept ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BG_COLOR)

    tb2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.2))
    tf2 = tb2.text_frame
    p = tf2.paragraphs[0]
    p.text = "WHAT IS A CITIZEN IN THIS ROOM?"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.font.name = "Segoe UI"
    p_sub = tf2.add_paragraph()
    p_sub.text = "Rights Without Responsibility Collapse a Community."
    p_sub.font.size = Pt(32)
    p_sub.font.bold = True
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.font.name = "Segoe UI"

    col_w = Inches(3.5)
    gap = Inches(0.4)
    start_x = Inches(1.0)
    y_pos = Inches(2.4)
    h_pos = Inches(4.2)

    cols = [
        ("Power & Voice", "You have direct input into how this room operates, how projects are assessed, and how disputes are handled. With voice comes the obligation to engage."),
        ("The Project Compact", "In Grade 9 Social Studies, collaborative simulations succeed or fail on individual accountability. Free-riders undermine the entire team's outcome."),
        ("Restitution vs. Punishment", "When norms break, how does our community respond? True civic justice focuses on repairing harm to the collective rather than passive compliance.")
    ]

    for i, (col_title, col_text) in enumerate(cols):
        cx = start_x + i * (col_w + gap)
        c_shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y_pos, col_w, h_pos)
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = CARD_BG
        c_shape.line.color.rgb = ACCENT_TEAL if i == 1 else ACCENT_GOLD
        c_shape.line.width = Pt(1)

        ctf = c_shape.text_frame
        ctf.word_wrap = True
        ctp = ctf.paragraphs[0]
        ctp.text = col_title
        ctp.font.size = Pt(18)
        ctp.font.bold = True
        ctp.font.color.rgb = ACCENT_GOLD if i != 1 else ACCENT_TEAL
        ctp.font.name = "Segoe UI"

        ctp_desc = ctf.add_paragraph()
        ctp_desc.space_before = Pt(14)
        ctp_desc.text = col_text
        ctp_desc.font.size = Pt(14)
        ctp_desc.font.color.rgb = TEXT_MUTED
        ctp_desc.font.name = "Segoe UI"

    # --- SLIDE 3: Student Action / Sticky Slide (STAYS ON SCREEN) ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, BG_COLOR)

    tb3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.2))
    tf3 = tb3.text_frame
    p = tf3.paragraphs[0]
    p.text = "STUDENT ACTION PROTOCOL"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.font.name = "Segoe UI"
    p_main = tf3.add_paragraph()
    p_main.text = "Step 1: 4-Corners Stand-Up  |  Step 2: 60-Second Civic Pitch"
    p_main.font.size = Pt(30)
    p_main.font.bold = True
    p_main.font.color.rgb = TEXT_WHITE
    p_main.font.name = "Segoe UI"

    # Left Card: 4-Corners
    card_l = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(5.4), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = ACCENT_TEAL
    card_l.line.width = Pt(1.5)

    ltf = card_l.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = "PHASE 1: 4-CORNERS STAND-UP DEBATE"
    lp.font.size = Pt(16)
    lp.font.bold = True
    lp.font.color.rgb = ACCENT_TEAL
    lp.font.name = "Segoe UI"

    l_items = [
        "Move physically to your corner (Strongly Agree, Agree, Disagree, Strongly Disagree):",
        "• Dilemma A: 'A group member who does zero work should receive a zero grade, even if the team project earns an A.'",
        "• Dilemma B: 'To protect class focus, phones should be placed in a shared organizer during collaborative work.'",
        "• Dilemma C: 'A student who disrupts a group owes direct work restitution to the team, not just an apology to the teacher.'",
        "Be ready to defend your corner with 1 piece of reasoning."
    ]
    for item in l_items:
        p = ltf.add_paragraph()
        p.space_before = Pt(8)
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE if "Dilemma" in item else TEXT_MUTED
        p.font.name = "Segoe UI"

    # Right Card: The Civic Pitch
    card_r = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = ACCENT_GOLD
    card_r.line.width = Pt(1.5)

    rtf = card_r.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = "PHASE 2: 60-SEC CIVIC PITCH & ROLE DRAFT"
    rp.font.size = Pt(16)
    rp.font.bold = True
    rp.font.color.rgb = ACCENT_GOLD
    rp.font.name = "Segoe UI"

    r_items = [
        "Turn to your assigned partner. Deliver your live verbal pitch (60 sec each):",
        "1. My Non-Negotiable: 'What is 1 essential condition you need from this class to do your best work?'",
        "2. My Project Role: 'What specific primary strength do you pledge to bring to group simulations (Lead Researcher, Visual Lead, Data/Logic Checker, Presenter, Project Manager)?'",
        "3. Cross-Examination: Partner must ask 1 challenge question before endorsing.",
        "4. Output: Post your pair's #1 non-negotiable on the front board."
    ]
    for item in r_items:
        p = rtf.add_paragraph()
        p.space_before = Pt(8)
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE if item.startswith(("1.", "2.", "3.", "4.")) else TEXT_MUTED
        p.font.name = "Segoe UI"

    cit9_path = os.path.join(out_dir, "CIT9_Day1_Civic_Contract.pptx")
    prs.save(cit9_path)
    print(f"Saved: {cit9_path}")

create_hl9_pptx()
create_cit9_pptx()
