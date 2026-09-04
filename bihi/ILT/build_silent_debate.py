"""Builds silent-debate-ilt.pptx — Silent Debate gallery lesson deck (16:9)."""
import os
from pptx import Presentation

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

DARK, TEXT, MUTED = RGBColor(0x2E, 0x38, 0x3E), RGBColor(0x3D, 0x44, 0x49), RGBColor(0x7F, 0x88, 0x8F)
ACCENT, TINT, HAIR = RGBColor(0xF5, 0x6A, 0x6A), RGBColor(0xFD, 0xEE, 0xEE), RGBColor(0xE5, 0xE7, 0xE9)
PANEL, GHOST, ONDARK, WHITE = RGBColor(0xF4, 0xF5, 0xF6), RGBColor(0xFA, 0xD8, 0xD8), RGBColor(0xAE, 0xB6, 0xBD), RGBColor(0xFF, 0xFF, 0xFF)
FONT, W = "Arial", 13.333

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(dark=False, hidden=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DARK if dark else WHITE
    if hidden:
        s._element.set("show", "0")
    return s


def tx(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "sb" in para:
            p.space_before = Pt(para["sb"])
        for t, o in para["runs"]:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name, f.size, f.bold, f.italic = FONT, Pt(o.get("size", 18)), o.get("bold", False), o.get("italic", False)
            f.color.rgb = o.get("color", TEXT)
    return tb


def box(s, x, y, w, h, fill=None, line=None, lw=1.0, kind=MSO_SHAPE.RECTANGLE, radius=None):
    sh = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill:
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb, sh.line.width = line, Pt(lw)
    if radius is not None:
        sh.adjustments[0] = radius
    sh.shadow.inherit = False
    return sh


def arrow(s, x1, y1, x2, y2, color=ACCENT, lw=2.0, dash=None, head=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb, c.line.width = color, Pt(lw)
    if dash:
        c.line.dash_style = dash
    if head:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle"}))
    c.shadow.inherit = False
    return c


def note(s, text):
    s.notes_slide.notes_text_frame.text = text


def rules_rows(s, rows, y0=1.5, rh=1.02, textw=8.6):
    for i, (head, gloss) in enumerate(rows):
        y = y0 + i * rh
        tx(s, 0.55, y - 0.03, 0.75, 0.6, [{"runs": [(str(i + 1), {"size": 26, "bold": True, "color": ACCENT})]}])
        tx(s, 1.45, y, textw, rh - 0.08, [
            {"runs": [(head, {"size": 19, "bold": True})]},
            {"runs": [(gloss, {"size": 13.5, "color": MUTED})], "sb": 3},
        ])
        if i < len(rows) - 1:
            box(s, 0.55, y + rh - 0.09, 10.4, 0.012, fill=HAIR)


def station(idx, statement, provocation, tag, hidden=False, letter=None):
    s = slide(hidden=hidden)
    tx(s, 0.62, 0.6, 6.0, 0.4, [{"runs": [(tag, {"size": 15, "bold": True, "color": ACCENT})]}])
    tx(s, 0.3, 0.7, 3.1, 4.7, [{"runs": [(letter or str(idx), {"size": 190, "bold": True, "color": ACCENT})]}],
       anchor=MSO_ANCHOR.MIDDLE)
    tx(s, 3.75, 1.55, 9.0, 2.8, [{"runs": [(statement, {"size": 42, "bold": True})]}])
    tx(s, 3.78, 4.55, 8.9, 0.8, [{"runs": [("\u2014  " + provocation, {"size": 17, "italic": True, "color": MUTED})]}])
    tx(s, 3.78, 6.55, 8.9, 0.4, [{"runs": [("Opening arguments, rebuttals and replies all live on this sheet.",
                                            {"size": 12.5, "color": MUTED})]}])
    return s


# ── 1 · Title ────────────────────────────────────────────────────────────────
s = slide(dark=True)
tx(s, 8.75, -0.75, 4.3, 4.6, [{"runs": [("\u201d", {"size": 300, "bold": True, "color": RGBColor(0x3A, 0x45, 0x52)})]}])
tx(s, 0.9, 0.9, 9.0, 0.4, [{"runs": [("INTEGRATED LEARNING TIME  ·  GRADE 9", {"size": 13, "bold": True, "color": ONDARK})]}])
tx(s, 0.86, 2.05, 11.6, 1.5, [{"runs": [("SILENT DEBATE", {"size": 72, "bold": True, "color": WHITE}),
                                        (".", {"size": 72, "bold": True, "color": ACCENT})]}])
tx(s, 0.9, 3.8, 10.5, 0.6, [{"runs": [("One hour. 28 people. Zero talking.", {"size": 26, "bold": True, "color": ACCENT})]}])
tx(s, 0.9, 4.6, 9.6, 0.9, [{"runs": [("Seven statements around the room. Groups of four. Your pencil is your voice.",
                                      {"size": 16, "color": ONDARK})]}])
note(s, "60-minute run sheet — 8 min setup (slides 1-5) · 21 min rotation (7 stations x 3 min) · 20 min verdict + "
        "60-second presentations · 10 min reflection.\nMATERIALS: 7 sheets of chart paper taped around the room at eye "
        "height, one pencil colour per student, a bell/timer.\nBEFORE CLASS: decide your 7 home stations and which "
        "group of 4 starts at each.")

# ── 2 · The pitch ────────────────────────────────────────────────────────────
s = slide()
tx(s, 0.5, 0.05, 4.2, 4.4, [{"runs": [("\u201c", {"size": 260, "bold": True, "color": GHOST})]}])
tx(s, 1.6, 2.5, 10.13, 2.4, [
    {"align": PP_ALIGN.CENTER, "runs": [("The loudest argument you\u2019ll ever have \u2014", {"size": 40, "bold": True})]},
    {"align": PP_ALIGN.CENTER, "runs": [("without saying a word.", {"size": 40, "bold": True, "italic": True, "color": ACCENT})]},
])
tx(s, 2.4, 5.25, 8.5, 0.8, [{"align": PP_ALIGN.CENTER,
                             "runs": [("No voices. No hands. One sheet of paper per statement \u2014 and your best reasons.",
                                       {"size": 17, "color": MUTED})]}])
note(s, "Hook — read it, pause, then move straight to the rules. Total silence starts when the first group reaches a sheet.")

# ── 3 · The rules ────────────────────────────────────────────────────────────
s = slide()
tx(s, 0.55, 0.5, 8.0, 0.7, [{"runs": [("The rules", {"size": 36, "bold": True})]}])
rules_rows(s, [
    ("Silence is absolute.", "From the first bell to the last \u2014 no whispering, no mouthing, no gestures."),
    ("Argue with ideas, not people.", "Attack the argument on the sheet \u2014 never the student who wrote it."),
    ("Every claim needs a because.", "An opinion without a reason is just noise. Write the because."),
    ("Disagree on purpose.", "Rebuttals are the point: \u201cThey say \u2026 , but \u2026 , because \u2026 .\u201d"),
    ("Your colour is your voice.", "Pick one pencil colour and keep it all hour \u2014 it shows who\u2019s winning."),
])
box(s, 10.4, 1.5, 2.38, 4.9, fill=TINT, kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tx(s, 10.68, 1.85, 1.85, 4.3, [
    {"runs": [("WHY SILENCE?", {"size": 14, "bold": True, "color": ACCENT})]},
    {"runs": [("Silence stops the loudest voice from winning \u2014 and pushes quiet students\u2019 ideas onto the page.",
               {"size": 13, "color": TEXT})], "sb": 10},
    {"runs": [("Every colour on a sheet is a different student. At the end, you\u2019ll see exactly who persuaded whom.",
               {"size": 13, "color": TEXT})], "sb": 10},
])
note(s, "Walk the rules fast (90 seconds). Hand out pencil colours as you go. Model one example on the board: "
        "\u2018Grades should be abolished BECAUSE \u2026\u2019 so they hear the shape of a claim + because.")

# ── 4 · How the room works ───────────────────────────────────────────────────
s = slide()
tx(s, 0.55, 0.5, 10.0, 0.7, [{"runs": [("How the room works", {"size": 36, "bold": True})]}])
tx(s, 0.55, 1.25, 12.2, 0.5, [{"runs": [("Groups of four. Each group claims one home sheet, then rotates on the bell until every sheet is visited.",
                                         {"size": 16, "color": MUTED})]}])
bw, gap, by = 1.15, 0.38, 2.0
x0 = (W - (7 * bw + 6 * gap)) / 2
for i in range(7):
    x = x0 + i * (bw + gap)
    b = box(s, x, by, bw, bw, fill=ACCENT, kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    tf = b.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(i + 1)
    r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(30), True, WHITE
    if i < 6:
        arrow(s, x + bw + 0.05, by + bw / 2, x + bw + gap - 0.05, by + bw / 2)
mid = by + bw
first_c, last_c = x0 + bw / 2, x0 + 6 * (bw + gap) + bw / 2
arrow(s, last_c, mid + bw / 2 + 0.05, last_c, 3.85, color=MUTED, lw=1.5, dash=MSO_LINE_DASH_STYLE.DASH)
arrow(s, last_c, 3.85, first_c, 3.85, color=MUTED, lw=1.5, dash=MSO_LINE_DASH_STYLE.DASH)
arrow(s, first_c, 3.85, first_c, mid + bw / 2 + 0.08, color=MUTED, lw=1.5, dash=MSO_LINE_DASH_STYLE.DASH, head=True)
tx(s, 5.2, 3.95, 3.0, 0.35, [{"align": PP_ALIGN.CENTER,
                              "runs": [("after station 7 \u2192 station 1", {"size": 12, "color": MUTED})]}])
for x, big, small in [(1.77, "7", "statements around the room"),
                      (5.17, "3 min", "on each sheet \u2014 move when the bell rings"),
                      (8.57, "21 min", "of rotation \u2014 every group visits every sheet")]:
    tx(s, x, 4.65, 3.0, 1.7, [
        {"align": PP_ALIGN.CENTER, "runs": [(big, {"size": 44, "bold": True, "color": ACCENT})]},
        {"align": PP_ALIGN.CENTER, "runs": [(small, {"size": 14, "color": TEXT})], "sb": 4},
    ])
tx(s, 1.5, 6.6, 10.3, 0.4, [{"align": PP_ALIGN.CENTER,
                             "runs": [("Home sheets get one last visit at the end \u2014 that\u2019s where the verdict happens.",
                                       {"size": 15, "color": MUTED})]}])
note(s, "Assign home stations BEFORE showing this. Rotation: bell every 3 minutes, groups move clockwise "
        "(1\u21922\u2192\u2026\u21927\u21921). 7 stops x 3 min = 21 minutes. Keep the timer visible.")

# ── 5 · Toolkit ──────────────────────────────────────────────────────────────
s = slide()
tx(s, 0.55, 0.5, 10.0, 0.7, [{"runs": [("Your argument toolkit", {"size": 36, "bold": True})]}])
steps = [
    ("CLAIM", "State it like it\u2019s already true. Short, certain, one sentence."),
    ("REASON", "Follow it with because. The reason is the argument."),
    ("EVIDENCE", "A fact, a number, something checkable \u2014 not just feelings."),
    ("EXAMPLE", "Make it concrete: this school, this week, your own life."),
]
sw, sgap, sy, sh = 2.55, 0.45, 1.65, 2.0
sx = (W - (4 * sw + 3 * sgap)) / 2
for i, (name, desc) in enumerate(steps):
    x = sx + i * (sw + sgap)
    b = box(s, x, sy, sw, sh, fill=PANEL, kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = name
    r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(20), True, ACCENT
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    r2 = p2.add_run()
    r2.text = desc
    r2.font.name, r2.font.size, r2.font.color.rgb = FONT, Pt(13), TEXT
    if i < 3:
        arrow(s, x + sw + 0.06, sy + sh / 2, x + sw + sgap - 0.06, sy + sh / 2)
box(s, 0, 4.25, W, 1.5, fill=TINT)
tx(s, 1.0, 4.25, 11.3, 1.5, [
    {"runs": [("THE STRONGEST MOVE", {"size": 13, "bold": True, "color": ACCENT})]},
    {"runs": [("\u201cThey say \u2026 , ", {"size": 22, "bold": True, "italic": True}),
              ("but", {"size": 22, "bold": True, "italic": True, "color": ACCENT}),
              (" \u2026 , ", {"size": 22, "bold": True, "italic": True}),
              ("because", {"size": 22, "bold": True, "italic": True, "color": ACCENT}),
              (" \u2026 .\u201d", {"size": 22, "bold": True, "italic": True})], "sb": 6},
])
tx(s, 1.0, 6.1, 11.3, 0.5, [{"runs": [("Concede before you counter \u2014 \u201cyou\u2019re right that \u2026 , but \u2026\u201d lands harder than plain disagreement.",
                                       {"size": 15, "color": MUTED})]}])
note(s, "2 minutes max. Read the chain aloud with one worked example: CLAIM \u2018Homework should be banned\u2019 \u2192 "
        "REASON \u2018because it eats the hours where teenagers actually recover\u2019 \u2192 EVIDENCE \u2192 EXAMPLE. "
        "The rebuttal formula is the slide to leave on screen during rotation.")

# ── 6-12 · Station slides ────────────────────────────────────────────────────
stations = [
    ("Grades should be abolished.", "If nobody got marks, who would still learn?"),
    ("School should start at 10 a.m.", "Teenage brains run on a different clock \u2014 so what are the bells for?"),
    ("AI tools should be allowed on every assignment.", "If the answer is findable, what is the assignment for?"),
    ("Video games should be an official school sport.", "Strategy, teamwork, reflexes \u2014 what\u2019s missing before it counts?"),
    ("Sixteen-year-olds should be able to vote.", "You can work and pay tax at sixteen \u2014 why not hold a ballot?"),
    ("Social media does more harm than good.", "The most documented generation in history \u2014 better, or worse?"),
    ("Homework should be banned.", "Six hours of school a day \u2014 is it enough time to learn?"),
]
for n, (st, pr) in enumerate(stations, 1):
    s = station(n, st, pr, f"STATION {n}")
    note(s, f"PRINT PAGE {n} \u2014 tape it to chart paper at Station {n}. Students open the debate here, then "
            f"rebut in other colours as sheets rotate.")

spares = [
    ("A", "Exams should be open-internet.", "You\u2019ll always have Google in real life \u2014 why not in the exam room?"),
    ("B", "Professional athletes are paid too much.", "Should a salary reflect skill, risk, or what people will pay to watch?"),
    ("C", "Every student must play on a school team.", "Is fitness a subject to teach, or a choice to protect?"),
]
for letter, st, pr in spares:
    s = station(letter, st, pr, "SPARE \u2014 SWAP FOR ANY STATION", hidden=True, letter=letter)
    note(s, "Hidden slide. Unhide and print if you want to swap out any of the 7 station statements.")

# ── 13 · The verdict ─────────────────────────────────────────────────────────
s = slide()
tx(s, 0.55, 0.5, 8.0, 0.7, [{"runs": [("The verdict", {"size": 36, "bold": True})]}])
tx(s, 0.55, 1.25, 12.0, 0.5, [{"runs": [("Twenty minutes left. Back to your home sheet.", {"size": 16, "color": MUTED})]}])
rules_rows(s, [
    ("Read everything on your sheet.", "Every colour, every rebuttal \u2014 especially the ones that trashed your opening argument."),
    ("Crown one winning argument.", "It can\u2019t be your own. Judge the writing, never the writer."),
    ("Name why it won.", "Logic, evidence, or a killer example \u2014 identify the move that beat you."),
    ("Nominate a speaker.", "Sixty seconds to present your winner. Everyone else: be ready to challenge it."),
], y0=2.0, rh=0.95, textw=11.3)
box(s, 0, 6.05, W, 1.0, fill=TINT)
tx(s, 1.0, 6.05, 11.3, 1.0, [{"align": PP_ALIGN.CENTER,
                              "runs": [("The judging test: would this argument still stand if it were read aloud to someone who wasn\u2019t in the room?",
                                        {"size": 16, "italic": True, "color": TEXT})]}], anchor=MSO_ANCHOR.MIDDLE)
note(s, "3 min silent reading + 2 min picking/prepping. Then 7 speakers x 60 seconds \u2248 8 min, with quick "
        "challenges between. Keep your own tally \u2014 it feeds the exit reflection.")

# ── 14 · Exit reflection ─────────────────────────────────────────────────────
s = slide(dark=True)
tx(s, 0.9, 0.85, 9.0, 0.4, [{"runs": [("BEFORE YOU GO", {"size": 13, "bold": True, "color": ACCENT})]}])
for i, (num, q) in enumerate([("01", "Which argument almost changed your mind \u2014 and what stopped it?"),
                              ("02", "What actually won today: logic, evidence, or an example?")]):
    y = 2.0 + i * 1.75
    tx(s, 0.9, y, 1.1, 0.9, [{"runs": [(num, {"size": 36, "bold": True, "color": ACCENT})]}])
    tx(s, 2.1, y + 0.05, 10.2, 1.3, [{"runs": [(q, {"size": 27, "bold": True, "color": WHITE})]}])
tx(s, 2.1, 5.9, 10.0, 0.5, [{"runs": [("Scrap paper. Two minutes. Hand it in on the way out.",
                                       {"size": 15, "color": ONDARK})]}])
note(s, "Exit ticket on scrap paper \u2014 2 minutes. Skim them tonight: the near-persuaded answers tell you which "
        "ILT project hooks the class (forensics, AI, sport analytics all appear in the ILT repository).")

out_pptx = os.path.join(os.path.dirname(os.path.abspath(__file__)), "silent-debate-ilt.pptx")
prs.save(out_pptx)
print("saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides to", out_pptx)

