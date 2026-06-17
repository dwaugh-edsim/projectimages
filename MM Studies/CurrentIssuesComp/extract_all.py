import os
import sys
from pathlib import Path

BASE = Path(r"E:\Antigravity\simroom\Github Repos\projectimages\MM Studies\CurrentIssuesComp")
OUT = BASE / "extracted_text.txt"

files = [
    "1. a Inquiry into Murdered and Missing Indigenous Women.pptx",
    "1. b The Inquiry into Missing and Murdered Aboriginal Women_ Answers.docx",
    "1. d Missing and Murdered Aboriginal Women in Canada.docx",
    "3. a Uncontacted Tribes.docx",
    "3. c Elon Musk brings Internet to Remote Tribes.docx",
    "4. a Change has started\u2013How have things been changing for the better in Mi\u2019kmaw Communities.docx",
    "5. a True Story_  2 Part Video.docx",
    "5. b High School Supplementary Resource Material for Treaty Education April 28.pdf",
    "5. f Pocahontas Myth.doc.docx",
    "5. g The Pocahontas Myth Answers.docx",
    "5. h The Pocahontas Myth Questions.docx",
    "7. a First Nations in the News Today.docx",
    "7. b In the News\u2026.docx",
    "7. c UN Declaration on the Rights of Indigenous Peoples.pdf",
    "Positives and Negatives Highlights from Pocahontas.docx",
    "Reel Injun Questions.doc.docx",
    "The Real Story of Pocahontas.docx",
]


def extract_docx(path):
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        text = para.text
        style = para.style.name if para.style else ""
        if style.startswith("Heading"):
            parts.append(f"[{style}] {text}")
        else:
            parts.append(text)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_texts = []
            for cell in row.cells:
                row_texts.append(cell.text)
            parts.append(" | ".join(row_texts))
    return "\n".join(parts)


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        row_texts.append(cell.text)
                    parts.append(" | ".join(row_texts))
    return "\n".join(parts)


def extract_pdf(path):
    import pdfplumber
    parts = []
    with pdfplumber.open(str(path)) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            parts.append(f"--- Page {page_num} ---")
            text = page.extract_text()
            if text:
                parts.append(text)
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    row_texts = [str(cell) if cell else "" for cell in row]
                    parts.append(" | ".join(row_texts))
    return "\n".join(parts)


with open(OUT, "w", encoding="utf-8") as out:
    for fname in files:
        fpath = BASE / fname
        if not fpath.exists():
            out.write(f"===FILE NOT FOUND=== {fname}\n\n")
            print(f"MISSING: {fname}", file=sys.stderr)
            continue

        print(f"Processing: {fname}", file=sys.stderr)
        ext = fpath.suffix.lower()
        try:
            if ext == ".docx":
                text = extract_docx(fpath)
            elif ext == ".pptx":
                text = extract_pdf(fpath) if False else extract_pptx(fpath)
            elif ext == ".pdf":
                text = extract_pdf(fpath)
            else:
                text = f"[Unsupported extension: {ext}]"

            out.write(f"==={fname}===\n")
            out.write(text)
            out.write(f"\n===END===\n\n")
        except Exception as e:
            out.write(f"==={fname}===\n")
            out.write(f"[ERROR: {e}]\n")
            out.write(f"===END===\n\n")
            print(f"ERROR on {fname}: {e}", file=sys.stderr)

print("Done.", file=sys.stderr)
