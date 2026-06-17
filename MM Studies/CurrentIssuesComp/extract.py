import os, json, sys

base = r"E:\Antigravity\simroom\Github Repos\projectimages\MM Studies\CurrentIssuesComp"

files = [
    "1. a Inquiry into Murdered and Missing Indigenous Women.pptx",
    "1. b The Inquiry into Missing and Murdered Aboriginal Women_ Answers.docx",
    "1. d Missing and Murdered Aboriginal Women in Canada.docx",
    "3. a Uncontacted Tribes.docx",
    "3. c Elon Musk brings Internet to Remote Tribes.docx",
    "4. a Change has started\u2013How have things been changing for the better in Mi\u2019kmaw Communities.docx",
    "5. a True Story_  2 Part Video.docx",
    "5. b High School Supplementary Resource Material for Treaty Education April 28.pdf",
    "5. f Pocahontas Myth.doc.docx",
    "5. g The Pocahontas Myth Answers.docx",
    "5. h The Pocahontas Myth Questions.docx",
    "7. a First Nations in the News Today.docx",
    "7. b In the News\u2026.docx",
    "7. c UN Declaration on the Rights of Indigenous Peoples.pdf",
    "Positives and Negatives Highlights from Pocahontas.docx",
    "Reel Injun Questions.doc.docx",
    "The Real Story of Pocahontas.docx",
]

results = {}

for fname in files:
    fpath = os.path.join(base, fname)
    if not os.path.exists(fpath):
        results[fname] = f"[FILE NOT FOUND: {fpath}]"
        continue
    try:
        if fname.endswith(".docx"):
            from docx import Document
            doc = Document(fpath)
            paragraphs = []
            for p in doc.paragraphs:
                text = p.text.strip()
                if text:
                    style = p.style.name if p.style else ""
                    paragraphs.append({"text": text, "style": style})
            # Also check tables
            tables = []
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                tables.append(rows)
            results[fname] = {"paragraphs": paragraphs, "tables": tables}
        elif fname.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(fpath)
            slides = []
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                    if shape.has_table:
                        tbl = shape.table
                        for row in tbl.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_texts.append(" | ".join(cells))
                slides.append(slide_texts)
            results[fname] = {"slides": slides}
        elif fname.endswith(".pdf"):
            import pdfplumber
            with pdfplumber.open(fpath) as pdf:
                pages = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        pages.append(text)
                results[fname] = {"pages": pages}
    except Exception as e:
        results[fname] = f"[ERROR: {str(e)}]"

output_path = os.path.join(base, "extracted_content.json")
with open(output_path, "w", encoding="utf-8") as f:
    json.dump(results, f, ensure_ascii=False, indent=2)

print(f"Extracted content from {len(results)} files to {output_path}")
for k, v in results.items():
    if isinstance(v, str) and v.startswith("["):
        print(f"  ISSUE: {k} -> {v}")
    elif isinstance(v, dict):
        if "paragraphs" in v:
            print(f"  OK: {k} -> {len(v['paragraphs'])} paragraphs, {len(v['tables'])} tables")
        elif "slides" in v:
            print(f"  OK: {k} -> {len(v['slides'])} slides")
        elif "pages" in v:
            print(f"  OK: {k} -> {len(v['pages'])} pages")
